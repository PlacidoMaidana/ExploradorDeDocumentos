import os
import tempfile
from utils.constants import EXT_VALIDAS
from core.file_utils import contar_archivos_validos

# Importar las nuevas bibliotecas
try:
    from docx import Document
    from pptx import Presentation
    from pdfminer.high_level import extract_text
    from pdfminer.layout import LAParams
    from PyPDF2 import PdfReader
    HAVE_NEW_DEPS = True
except ImportError:
    HAVE_NEW_DEPS = False

def extraer_texto_docx(ruta):
    """Extrae texto de archivos Word (.docx)"""
    try:
        doc = Document(ruta)
        texto_completo = []
        
        # Extraer texto de párrafos
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texto_completo.append(paragraph.text)
        
        # Extraer texto de tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texto_completo.append(cell.text)
        
        return '\n'.join(texto_completo)
    except Exception as e:
        return f"[ERROR] No se pudo procesar el archivo Word: {str(e)}"

def extraer_texto_pptx(ruta):
    """Extrae texto de archivos PowerPoint (.pptx)"""
    try:
        prs = Presentation(ruta)
        texto_completo = []
        
        for i, slide in enumerate(prs.slides):
            texto_completo.append(f"--- Diapositiva {i+1} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texto_completo.append(shape.text)
        
        return '\n'.join(texto_completo)
    except Exception as e:
        return f"[ERROR] No se pudo procesar el archivo PowerPoint: {str(e)}"

def extraer_texto_pdf(ruta):
    """Extrae texto y comentarios de archivos PDF con mejor manejo de errores"""
    try:
        texto_completo = []
        
        # Verificar si el archivo existe y es legible
        if not os.path.exists(ruta):
            return "[ERROR] El archivo PDF no existe"
        
        # Intentar extraer texto principal con pdfminer
        try:
            laparams = LAParams()
            texto = extract_text(ruta, laparams=laparams)
            if texto and texto.strip():
                texto_completo.append("=== TEXTO PRINCIPAL ===")
                texto_completo.append(texto.strip())
            else:
                texto_completo.append("=== TEXTO PRINCIPAL ===")
                texto_completo.append("No se pudo extraer texto del PDF. El archivo puede estar encriptado, ser una imagen o estar corrupto.")
        except Exception as e:
            texto_completo.append(f"[ERROR pdfminer] No se pudo extraer texto: {str(e)}")
        
        # Intentar extraer metadatos y estructura con PyPDF2
        try:
            with open(ruta, 'rb') as file:
                pdf_reader = PdfReader(file)
                
                # Información básica del PDF
                info = pdf_reader.metadata
                if info:
                    texto_completo.append("\n=== METADATOS DEL PDF ===")
                    for key, value in info.items():
                        texto_completo.append(f"{key}: {value}")
                
                # Verificar si está encriptado
                if pdf_reader.is_encrypted:
                    texto_completo.append("\n[ADVERTENCIA] El PDF está encriptado/protegido")
                    # Intentar desencriptar con contraseña vacía (para algunos PDFs)
                    try:
                        if pdf_reader.decrypt(''):
                            texto_completo.append("[INFO] PDF desencriptado con contraseña vacía")
                        else:
                            texto_completo.append("[INFO] No se pudo desencriptar el PDF")
                    except:
                        texto_completo.append("[INFO] El PDF requiere contraseña para desencriptar")
                
                # Intentar extraer texto página por página
                texto_pypdf = []
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            texto_pypdf.append(f"--- Página {page_num} ---")
                            texto_pypdf.append(page_text.strip())
                    except Exception as page_error:
                        texto_pypdf.append(f"[ERROR Página {page_num}] {str(page_error)}")
                
                if texto_pypdf:
                    texto_completo.append("\n=== TEXTO EXTRAÍDO CON PyPDF2 ===")
                    texto_completo.extend(texto_pypdf)
                
                # Intentar extraer comentarios/anotaciones
                try:
                    comentarios = []
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        if '/Annots' in page:
                            for annot in page['/Annots']:
                                annot_obj = annot.get_object()
                                if '/Contents' in annot_obj:
                                    comentario = annot_obj['/Contents']
                                    if comentario:
                                        comentarios.append(f"Página {page_num}: {comentario}")
                    
                    if comentarios:
                        texto_completo.append("\n=== COMENTARIOS/ANOTACIONES ===")
                        texto_completo.extend(comentarios)
                except Exception as e:
                    texto_completo.append(f"\n[INFO] No se pudieron extraer comentarios: {str(e)}")
                        
        except Exception as e:
            texto_completo.append(f"\n[ERROR PyPDF2] {str(e)}")
        
        # Si no se pudo extraer ningún texto, proporcionar información de diagnóstico
        if len(texto_completo) <= 2:  # Solo tiene encabezados básicos
            texto_completo.append("\n=== INFORMACIÓN DE DIAGNÓSTICO ===")
            texto_completo.append("El PDF podría ser:")
            texto_completo.append("- Un documento escaneado (imagen) sin texto extraíble")
            texto_completo.append("- Un PDF encriptado/protegido contra extracción")
            texto_completo.append("- Un PDF corrupto o dañado")
            texto_completo.append("- Un documento que requiere OCR")
            
            # Información del archivo
            try:
                file_size = os.path.getsize(ruta)
                texto_completo.append(f"\nTamaño del archivo: {file_size} bytes")
            except:
                pass
        
        return '\n'.join(texto_completo)
        
    except Exception as e:
        return f"[ERROR CRÍTICO] No se pudo procesar el archivo PDF: {str(e)}"

def leer_archivo_especial(ruta, ext):
    """Lee archivos especiales (Word, PowerPoint, PDF) según su extensión"""
    if not HAVE_NEW_DEPS:
        return f"[ERROR] Bibliotecas necesarias no instaladas. Ejecuta: pip install python-docx python-pptx pdfminer.six PyPDF2"
    
    try:
        if ext == '.docx':
            return extraer_texto_docx(ruta)
        elif ext == '.pptx':
            return extraer_texto_pptx(ruta)
        elif ext == '.pdf':
            resultado = extraer_texto_pdf(ruta)
            # Si no se pudo extraer texto, intentar con PyMuPDF
            if "No se pudo extraer texto" in resultado or "encriptado" in resultado.lower():
                alternativa = procesar_pdf_con_ocm(ruta)
                if "EXTRAÍDO" in alternativa:
                    resultado += f"\n\n{alternativa}"
            return resultado
        else:
            return f"[ERROR] Formato no soportado: {ext}"
    except Exception as e:
        return f"[ERROR] No se pudo procesar el archivo: {str(e)}"

def mostrar_contenido_archivos(ruta, progress_callback=None, excluir=None, incluir=None):
    """Muestra contenido de todos los archivos válidos en el árbol"""
    excluir = set(excluir or [])
    incluir = set(incluir or [])
    salida = ''
    total_files = contar_archivos_validos(ruta)
    processed = 0

    for root, dirs, files in os.walk(ruta):
        # Filtrar directorios
        dirs_filtrados = []
        for d in dirs:
            ruta_relativa = os.path.relpath(os.path.join(root, d), ruta)
            if ruta_relativa in excluir:
                continue
            if incluir and ruta_relativa not in incluir:
                continue
            dirs_filtrados.append(d)
        dirs[:] = dirs_filtrados

        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in EXT_VALIDAS:
                full_path = os.path.join(root, file)
                ruta_relativa = os.path.relpath(full_path, ruta)
                
                try:
                    # Para archivos especiales (Word, PowerPoint, PDF)
                    if ext in ['.docx', '.pptx', '.pdf']:
                        contenido = leer_archivo_especial(full_path, ext)
                    else:
                        # Para archivos de texto (código anterior)
                        try:
                            with open(full_path, 'r', encoding='utf-8') as f:
                                contenido = f.read()
                        except UnicodeDecodeError:
                            try:
                                with open(full_path, 'r', encoding='latin-1') as f:
                                    contenido = f.read()
                            except Exception as e:
                                contenido = f"[ERROR] No se pudo leer '{ruta_relativa}': {str(e)}"
                    
                    salida += f"\n{'='*80}\n📄 {ruta_relativa}\n{'='*80}\n\n{contenido.strip()}\n\n"
                    
                except Exception as e:
                    salida += f"\n[ERROR] No se pudo procesar '{ruta_relativa}': {str(e)}\n\n"

                processed += 1
                if progress_callback:
                    progress_callback(processed, total_files)

    return salida

def mostrar_contenido_raiz(ruta, progress_callback=None):
    """Muestra contenido solo de archivos en la raíz del directorio"""
    salida = ''
    archivos_validos = []

    for file in os.listdir(ruta):
        full_path = os.path.join(ruta, file)
        if os.path.isfile(full_path):
            ext = os.path.splitext(file)[1].lower()
            if ext in EXT_VALIDAS:
                archivos_validos.append((file, full_path))

    total_files = len(archivos_validos)

    for i, (file, full_path) in enumerate(archivos_validos):
        try:
            ext = os.path.splitext(file)[1].lower()
            
            # Para archivos especiales (Word, PowerPoint, PDF)
            if ext in ['.docx', '.pptx', '.pdf']:
                contenido = leer_archivo_especial(full_path, ext)
            else:
                # Para archivos de texto (código anterior)
                try:
                    with open(full_path, 'r', encoding='utf-8') as f:
                        contenido = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(full_path, 'r', encoding='latin-1') as f:
                            contenido = f.read()
                    except Exception as e:
                        contenido = f"[ERROR] No se pudo leer '{file}': {str(e)}"
            
            salida += f"\n{'='*80}\n📄 {file}\n{'='*80}\n\n{contenido.strip()}\n\n"
            
        except Exception as e:
            salida += f"\n[ERROR] No se pudo procesar '{file}': {str(e)}\n\n"

        if progress_callback:
            progress_callback(i + 1, total_files)

    return salida


def procesar_pdf_con_ocm(ruta):
    """Intenta procesar PDFs problemáticos usando PyMuPDF (fitz) como alternativa"""
    try:
        import fitz  # PyMuPDF
        texto_completo = ["=== TEXTO EXTRAÍDO CON PyMuPDF ==="]
        
        doc = fitz.open(ruta)
        for page_num in range(len(doc)):
            page = doc[page_num]
            texto = page.get_text()
            if texto.strip():
                texto_completo.append(f"--- Página {page_num + 1} ---")
                texto_completo.append(texto.strip())
        
        doc.close()
        return '\n'.join(texto_completo)
    except ImportError:
        return "[INFO] PyMuPDF no está instalado. Ejecuta: pip install PyMuPDF"
    except Exception as e:
        return f"[ERROR PyMuPDF] {str(e)}"
    
def verificar_tipo_archivo(ruta):
    """Verifica el tipo real del archivo"""
    try:
        import magic
        file_type = magic.from_file(ruta, mime=True)
        return file_type
    except ImportError:
        return "[INFO] python-magic no instalado"
    except Exception as e:
        return f"[ERROR] {str(e)}"    